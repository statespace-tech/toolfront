import json
import logging
import xml.etree.ElementTree as ET
from abc import ABC
from pathlib import Path
from typing import Any
from urllib.parse import ParseResult, urlparse

from pydantic import BaseModel, Field, field_validator

from toolfront.models.database import SearchMode
from toolfront.types import ConnectionResult, DocumentType
from toolfront.utils import search_items

logger = logging.getLogger("toolfront")


class LibraryError(Exception):
    """Exception for library-related errors."""

    pass


class Library(BaseModel, ABC):
    """Abstract base class for library."""

    url: ParseResult = Field(description="URL of the library")

    @field_validator("url", mode="before")
    def validate_url(cls, v: Any) -> ParseResult:
        if isinstance(v, str):
            v = urlparse(v)
        return v

    def _parse_pagination(self, pagination: int | float) -> tuple[float | None, int | None]:
        """Parse and validate pagination parameter.

        Args:
            pagination: Page/section number (1+ int) or percentile (0-1 exclusive float)

        Returns:
            Tuple of (page_percentile, page_number) where exactly one is None

        Rules:
            - Float 0 < pagination < 1: percentile
            - Int pagination >= 1: page number
            - Invalid values default to page 1
        """
        if isinstance(pagination, float):
            if 0 < pagination < 1:
                return pagination, None
            else:
                # Invalid percentile or edge case, default to page 1
                return None, 1
        else:
            # Integer case
            if pagination >= 1:
                return None, pagination
            else:
                # Invalid page number, default to page 1
                return None, 1

    def _get_target_page(self, page_percentile: float | None, page_number: int | None, total_pages: int) -> int:
        """Get the target page number, defaulting to page 1 if no pagination specified."""
        if page_percentile is not None:
            return max(1, min(total_pages, int(page_percentile * total_pages) + 1))
        elif page_number is not None:
            return max(1, min(total_pages, page_number))
        else:
            return 1  # Always default to page 1

    async def test_connection(self) -> ConnectionResult:
        """Test the connection to the library."""
        return ConnectionResult(connected=True, message="Library connection successful")

    async def get_documents(self) -> list[str]:
        """Get all documents in the library recursively."""
        path = Path(self.url.path)
        if not path.exists():
            return []

        try:
            supported_extensions = DocumentType.get_supported_extensions()
            return [str(p.relative_to(path)) for p in path.rglob("*.*") if p.suffix.lower() in supported_extensions]
        except (PermissionError, OSError) as e:
            logger.warning(f"Error accessing {path}: {e}")
            return []

    async def search_documents(self, pattern: str, mode: SearchMode = SearchMode.REGEX, limit: int = 10) -> list[str]:
        """Search for documents in the library."""
        files = await self.get_documents()
        return search_items(files, pattern, mode, limit)

    def _read_docx(self, path: str) -> str:
        """Read DOCX document."""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(path)

            # Read all paragraphs
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"

            return text
        except ImportError:
            return "DOCX support requires python-docx library"
        except Exception as e:
            return f"Error reading DOCX file: {str(e)}"

    def _read_excel(self, path: str, pagination: int | float = 0) -> str:
        """Read a single sheet from Excel document."""
        try:
            import pandas as pd

            # Parse pagination parameter
            page_percentile, page_number = self._parse_pagination(pagination)

            excel_data = pd.read_excel(path, sheet_name=None)
            sheet_names = list(excel_data.keys())
            total_sheets = len(sheet_names)

            target_sheet_idx = self._get_target_page(page_percentile, page_number, total_sheets) - 1
            target_sheet_name = sheet_names[target_sheet_idx]
            df = excel_data[target_sheet_name]

            result = f"Sheet {target_sheet_idx + 1} of {total_sheets}: {target_sheet_name}\n"
            result += f"Shape: {df.shape}\n"
            result += f"Columns: {list(df.columns)}\n"
            result += df.to_string() + "\n"

            return result
        except ImportError:
            return "Excel support requires pandas and openpyxl libraries"
        except Exception as e:
            return f"Error reading Excel file: {str(e)}"

    def _read_json(self, path: str) -> str:
        """Read JSON document."""
        try:
            with Path(path).open("r", encoding="utf-8") as file:
                data = json.load(file)
                return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError as e:
            return f"Error parsing JSON: {str(e)}"
        except Exception as e:
            return f"Error reading JSON file: {str(e)}"

    def _read_markdown(self, path: str) -> str:
        """Read Markdown document."""
        try:
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()
        except Exception as e:
            return f"Error reading Markdown file: {str(e)}"

    def _read_pdf(self, path: str, pagination: int | float = 0) -> str:
        """Read a single page from PDF document."""
        try:
            from pypdf import PdfReader

            # Parse pagination parameter
            page_percentile, page_number = self._parse_pagination(pagination)

            reader = PdfReader(path)
            total_pages = len(reader.pages)

            target_page_idx = self._get_target_page(page_percentile, page_number, total_pages) - 1

            text = f"Page {target_page_idx + 1} of {total_pages}:\n\n"
            text += reader.pages[target_page_idx].extract_text()

            return text
        except ImportError:
            return "PDF support requires pypdf library"
        except Exception as e:
            return f"Error reading PDF file: {str(e)}"

    def _read_powerpoint(self, path: str, pagination: int | float = 0) -> str:
        """Read a single slide from PowerPoint document."""
        try:
            from pptx import Presentation

            # Parse pagination parameter
            page_percentile, page_number = self._parse_pagination(pagination)

            prs = Presentation(path)
            total_slides = len(prs.slides)

            target_slide_idx = self._get_target_page(page_percentile, page_number, total_slides) - 1

            slide = prs.slides[target_slide_idx]
            text = f"Slide {target_slide_idx + 1} of {total_slides}:\n\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

            return text
        except ImportError:
            return "PowerPoint support requires python-pptx library"
        except Exception as e:
            return f"Error reading PowerPoint file: {str(e)}"

    def _read_rtf(self, path: str) -> str:
        """Read RTF document."""
        try:
            from striprtf.striprtf import rtf_to_text

            with Path(path).open("r", encoding="utf-8") as file:
                rtf_content = file.read()
                text = rtf_to_text(rtf_content)
                return text
        except ImportError:
            # Fallback to raw RTF if striprtf not available
            with Path(path).open("r", encoding="utf-8") as file:
                return f"RTF content (requires striprtf library for text extraction):\n{file.read()}"
        except Exception as e:
            return f"Error reading RTF file: {str(e)}"

    def _read_txt(self, path: str) -> str:
        """Read TXT document."""
        try:
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()
        except Exception as e:
            return f"Error reading TXT file: {str(e)}"

    def _read_xml(self, path: str) -> str:
        """Read XML document."""
        try:
            tree = ET.parse(path)
            root = tree.getroot()

            def element_to_string(elem, level=0):
                indent = "  " * level
                result = f"{indent}Element: {elem.tag}\n"
                if elem.attrib:
                    result += f"{indent}Attributes: {elem.attrib}\n"
                if elem.text and elem.text.strip():
                    result += f"{indent}Text: {elem.text.strip()}\n"
                for child in elem:
                    result += element_to_string(child, level + 1)
                return result

            return element_to_string(root)
        except Exception:
            # Fallback to plain text if parsing fails
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()

    def _read_yaml(self, path: str) -> str:
        """Read YAML document."""
        try:
            import yaml

            with Path(path).open("r", encoding="utf-8") as file:
                data = yaml.safe_load(file)
                return yaml.dump(data, indent=2, default_flow_style=False)
        except ImportError:
            # Fallback to plain text if PyYAML not available
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()
        except Exception as e:
            return f"Error parsing YAML: {str(e)}"

    async def read_document(self, document_path: str, document_type: DocumentType, pagination: int | float = 0) -> str:
        """Read the file using appropriate method based on file extension.

        Args:
            document_type: Type of the document to read.
            document_path: Path to the document.
            pagination: Page/section number (1+ int) or percentile (0-1 exclusive float) to read.
                        Only used for paginated documents (PDF, PPTX, XLSX). Ignored for others.
        """

        full_path = Path(self.url.path) / document_path

        match document_type:
            case DocumentType.DOCX:
                return self._read_docx(full_path)
            case DocumentType.XLSX | DocumentType.XLS:
                return self._read_excel(full_path, pagination)
            case DocumentType.JSON:
                return self._read_json(full_path)
            case DocumentType.MD:
                return self._read_markdown(full_path)
            case DocumentType.PDF:
                return self._read_pdf(full_path, pagination)
            case DocumentType.PPTX:
                return self._read_powerpoint(full_path, pagination)
            case DocumentType.RTF:
                return self._read_rtf(full_path)
            case DocumentType.TXT:
                return self._read_txt(full_path)
            case DocumentType.XML:
                return self._read_xml(full_path)
            case DocumentType.YAML | DocumentType.YML:
                return self._read_yaml(full_path)
            case _:
                return f"Unsupported document type: {document_type}"
